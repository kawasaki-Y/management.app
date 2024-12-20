import streamlit as st
import pandas as pd
import altair as alt
from database import Database
from datetime import date
from pptx import Presentation
from pptx.util import Inches
import os

def save_chart_as_image(chart, filename):
    """Altairのグラフを画像として保存"""
    chart.save(filename, scale_factor=2)

def generate_pptx_with_charts(comparison_chart, monthly_chart, pie_chart):
    """PowerPointファイルを生成"""
    ppt = Presentation()

    # スライド1: 予実比較
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    title = slide.shapes.title
    title.text = "予実比較"
    save_chart_as_image(comparison_chart, "comparison_chart.png")
    slide.shapes.add_picture("comparison_chart.png", Inches(1), Inches(1.5), width=Inches(6))
    os.remove("comparison_chart.png")

    # スライド2: 月毎売上
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    title = slide.shapes.title
    title.text = "月毎売上"
    save_chart_as_image(monthly_chart, "monthly_chart.png")
    slide.shapes.add_picture("monthly_chart.png", Inches(1), Inches(1.5), width=Inches(6))
    os.remove("monthly_chart.png")

    # スライド3: タグ別売上比率
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    title = slide.shapes.title
    title.text = "タグ別売上比率"
    save_chart_as_image(pie_chart, "pie_chart.png")
    slide.shapes.add_picture("pie_chart.png", Inches(1), Inches(1.5), width=Inches(6))
    os.remove("pie_chart.png")

    # PowerPointファイルを保存
    ppt.save("sales_charts.pptx")
    return "sales_charts.pptx"

def sales_management_page():
    st.header("売上管理")

    # 登録済み目標売上を取得
    target_data = Database.fetch_data("SELECT amount FROM target_revenue")
    if not target_data:
        st.warning("目標売上が登録されていません。'タグと目標売上の登録' ページで登録してください。")
        return
    target_revenue = target_data[0][0]  # 千円単位の目標売上
    st.write(f"年間目標売上高: {target_revenue:,} 千円")

    # タグリストを取得
    tags_data = Database.fetch_data("SELECT tag_name FROM tags")
    tags = [row[0] for row in tags_data]
    if not tags:
        st.warning("タグが登録されていません。'タグと目標売上の登録' ページで登録してください。")
        return

    # 売上登録フォーム
    st.subheader("売上登録")
    if "form_submitted" not in st.session_state:
        st.session_state["form_submitted"] = False

    if not st.session_state["form_submitted"]:
        with st.form("sales_form"):
            project = st.text_input("案件名")
            selected_tag = st.selectbox("タグを選択", options=tags)
            revenue = st.number_input("売上高（千円単位）", min_value=0, step=100)
            date_input = st.date_input("日付", value=date.today())
            submitted_sales = st.form_submit_button("登録")
            if submitted_sales:
                Database.execute_query(
                    "INSERT INTO sales (project, tag, revenue, date) VALUES (?, ?, ?, ?)",
                    (project, selected_tag, revenue, str(date_input))
                )
                st.session_state["form_submitted"] = True
                st.success(f"案件 '{project}' を登録しました！")

    if st.session_state["form_submitted"]:
        if st.button("新規登録を行う"):
            st.session_state["form_submitted"] = False

    # 売上データの取得
    sales_data = Database.fetch_data("SELECT id, project, tag, revenue, date FROM sales")
    df = pd.DataFrame(sales_data, columns=["ID", "案件名", "タグ", "売上高（千円）", "日付"])
    if not df.empty:
        st.subheader("売上データ")

        # データの表示・修正・削除
        for i, row in df.iterrows():
            col1, col2, col3 = st.columns([4, 1, 1])
            with col1:
                st.text(f"案件名: {row['案件名']}, タグ: {row['タグ']}, 売上高: {row['売上高（千円）']:,} 千円, 日付: {row['日付']}")
            with col2:
                if st.button("修正", key=f"edit_{row['ID']}"):
                    new_project = st.text_input("案件名", value=row["案件名"], key=f"project_{row['ID']}")
                    new_tag = st.selectbox("タグ", options=tags, index=tags.index(row["タグ"]), key=f"tag_{row['ID']}")
                    new_revenue = st.number_input("売上高（千円単位）", min_value=0, value=row["売上高（千円）"], step=100, key=f"revenue_{row['ID']}")
                    new_date = st.date_input("日付", value=pd.to_datetime(row["日付"]), key=f"date_{row['ID']}")
                    if st.button("更新", key=f"update_{row['ID']}"):
                        Database.execute_query(
                            "UPDATE sales SET project = ?, tag = ?, revenue = ?, date = ? WHERE id = ?",
                            (new_project, new_tag, new_revenue, str(new_date), row["ID"])
                        )
                        st.success(f"案件 '{new_project}' を更新しました！")
                        st.experimental_rerun()
            with col3:
                if st.button("削除", key=f"delete_{row['ID']}"):
                    Database.execute_query("DELETE FROM sales WHERE id = ?", (row["ID"],))
                    st.success(f"案件 '{row['案件名']}' を削除しました！")
                    st.experimental_rerun()

        # グラフの作成
        st.subheader("売上グラフ")
        total_revenue = df["売上高（千円）"].sum()
        comparison_data = pd.DataFrame({
            "項目": ["目標売上高", "実績売上高"],
            "金額（千円）": [target_revenue, total_revenue]
        })
        comparison_chart = alt.Chart(comparison_data).mark_bar(color="#ff7f7f").encode(
            x="項目:N",
            y="金額（千円）:Q"
        )
        st.altair_chart(comparison_chart, use_container_width=True)

        # 月毎売上（折れ線グラフ）
        df["月"] = pd.to_datetime(df["日付"]).dt.to_period("M")
        monthly_sales = df.groupby("月")["売上高（千円）"].sum().reset_index()
        monthly_chart = alt.Chart(monthly_sales).mark_line(color="#ff7f7f").encode(
            x="月:T",
            y="売上高（千円）:Q"
        )
        st.altair_chart(monthly_chart, use_container_width=True)

        # タグ別売上比率（円グラフ）
        tag_sales = df.groupby("タグ")["売上高（千円）"].sum().reset_index()
        pie_chart = alt.Chart(tag_sales).mark_arc(innerRadius=50).encode(
            theta="売上高（千円）:Q",
            color="タグ:N",
            tooltip=["タグ", "売上高（千円）"]
        )
        st.altair_chart(pie_chart, use_container_width=True)

        # グラフをPowerPointにダウンロード
        st.subheader("グラフのエクスポート")
        if st.button("PowerPointファイルをダウンロード"):
            pptx_file = generate_pptx_with_charts(comparison_chart, monthly_chart, pie_chart)
            with open(pptx_file, "rb") as file:
                st.download_button(
                    label="PowerPointをダウンロード",
                    data=file,
                    file_name=pptx_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            os.remove(pptx_file)

        # SQLファイルダウンロード
        st.subheader("データエクスポート")
        if st.button("SQLファイルをダウンロード"):
            sql_file_path = "sales_data.sql"
            with open(sql_file_path, "w") as f:
                conn = Database.get_connection()
                for line in conn.iterdump():
                    f.write(f"{line}\n")
                conn.close()
            with open(sql_file_path, "rb") as file:
                st.download_button(
                    label="SQLファイルをダウンロード",
                    data=file,
                    file_name=sql_file_path,
                    mime="application/sql"
                )
            os.remove(sql_file_path)
